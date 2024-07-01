from pptx import Presentation
from transformers import pipeline

# Initialize the summarization pipeline with specific model and device
model_name = "sshleifer/distilbart-cnn-12-6"
summarizer = pipeline("summarization", model=model_name, device=0)  # Use device=0 for GPU, device=-1 for CPU

# Input text
text = """
Create a project in the Google Cloud Console. 
Enable the Google Slides API for your project. 
Set up OAuth 2.0 credentials for your application. 
Create a project in the Google Cloud Console. 
Enable the Google Slides API for your project.
Set up OAuth 2.0 credentials for your application.
"""

# Generate summary or structure
summary = summarizer(text, max_length=38, min_length=30, do_sample=False)[0]['summary_text']  # Adjust max_length accordingly
slides = summary.split(". ")  # Splitting summary into sentences for slide content

# Create a PowerPoint presentation
prs = Presentation()

for idx, slide_content in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.shapes.placeholders[1]
    title.text = f"Slide {idx + 1}"  # Dynamic slide titles
    content.text = slide_content

# Save the presentation
prs.save('output.pptx')
