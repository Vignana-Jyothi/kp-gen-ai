from pptx import Presentation
from pptx.util import Inches
import re

# Function to parse the file content
def parse_file(file_path):
    with open(file_path, 'r') as file:
        content = file.read()
    
    # Regular expressions to match slide titles and contents
    slide_pattern = re.compile(r'### Slide \d+: (.+?)\n\*\*Title:\*\* (.+?)\n\n\*\*Content:\*\*\n- (.+?)(?=###|$)', re.DOTALL)
    slides = slide_pattern.findall(content)

    parsed_slides = []
    for slide in slides:
        slide_number, title, body = slide
        body = body.strip().split('\n- ')
        parsed_slides.append((title, body))
        print("------")
        print(title)
        print(body)
    
    return parsed_slides

# Function to generate PowerPoint from parsed content
def generate_ppt(parsed_slides, output_path):
    prs = Presentation()
    for title, content in parsed_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        content_shape = slide.placeholders[1]
        
        # Join content list into a single string with newlines
        content_text = '\n'.join(content)
        content_shape.text = content_text

    prs.save(output_path)

# File path and output path
file_path = 'slides_content.txt'
output_path = 'output_presentation.pptx'

# Parse the file and generate the PowerPoint
parsed_slides = parse_file(file_path)
generate_ppt(parsed_slides, output_path)
